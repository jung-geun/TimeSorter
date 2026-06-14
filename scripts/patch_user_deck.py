#!/usr/bin/env python
"""사용자 제공 PPT의 실험 섹션만 최신화 (앞부분 데이터 설명·이미지 보존).

- 슬라이드 1~12: 손대지 않음 (사용자 커스텀 DATA SOURCE 이미지 포함)
- 슬라이드 13~15: 차트 이미지 교체(n=150) + 캡션 수정
- 슬라이드 17~18: 표 수치 n=150 갱신
- 슬라이드 19(결론): 수치 갱신
- 신규: 4B vs 9B, 9B 4-way, 학습 지표 슬라이드를 결론 앞에 삽입
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Pt, Inches

SRC = "/tmp/user_deck.pptx"
OUT = "/mnt/hdd/WD_8TB/code/TimeSorter/presentation/TimeSorter_발표자료_v2.pptx"
ASSETS = Path("/mnt/hdd/WD_8TB/code/TimeSorter/assets")

NAVY = RGBColor(0x1F, 0x2D, 0x3D); BLUE = RGBColor(0x4C, 0x72, 0xB0)
ORANGE = RGBColor(0xDD, 0x84, 0x52); GREEN = RGBColor(0x55, 0xA8, 0x68)
RED = RGBColor(0xC0, 0x39, 0x2B); GRAY = RGBColor(0x5A, 0x5A, 0x5A)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)
FONT = "Noto Sans CJK KR"

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
slides = list(prs.slides)


# ── 헬퍼 ──────────────────────────────────────────────────────────────────────
def _set_font(run, size=18, bold=False, color=DARK, font=FONT, italic=False):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font


def add_box(slide, x, y, w, h, fill=None, line=None, line_w=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = line_w or Pt(1)
    return shp


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for i, (text, kw) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = kw.pop("align", align)
        if "space_after" in kw: p.space_after = Pt(kw.pop("space_after"))
        if "space_before" in kw: p.space_before = Pt(kw.pop("space_before"))
        p.level = kw.pop("level", 0)
        bullet = kw.pop("bullet", False)
        run = p.add_run(); run.text = ("•  " if bullet else "") + text
        _set_font(run, **kw)
    return tb


def add_table(slide, x, y, w, h, data, col_widths=None, header_fill=BLUE, font_size=13, header_size=13):
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        for ci, cw in enumerate(col_widths): gt.columns[ci].width = cw
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.add_run(); run.text = str(val)
            if ri == 0:
                _set_font(run, size=header_size, bold=True, color=WHITE)
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                _set_font(run, size=font_size, color=DARK)
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
    return gt


def add_image_fit(slide, path, x, y, max_w, max_h):
    from PIL import Image
    iw, ih = Image.open(path).size
    r = min(max_w / iw, max_h / ih); w, h = int(iw * r), int(ih * r)
    slide.shapes.add_picture(str(path), x + (max_w - w) // 2, y + (max_h - h) // 2, w, h)


def title_bar(slide, kicker, title):
    add_box(slide, 0, 0, SW, Inches(1.15), fill=NAVY)
    add_box(slide, 0, Inches(1.15), SW, Pt(3), fill=ORANGE)
    if kicker:
        add_text(slide, Inches(0.5), Inches(0.14), Inches(12), Inches(0.35),
                 [(kicker, dict(size=12, bold=True, color=ORANGE))])
    add_text(slide, Inches(0.5), Inches(0.42), Inches(12.3), Inches(0.65),
             [(title, dict(size=26, bold=True, color=WHITE))])


def new_slide(kicker, title):
    s = prs.slides.add_slide(BLANK); title_bar(s, kicker, title); return s


# ── 이미지/텍스트 편집 헬퍼 ───────────────────────────────────────────────────
def replace_pic(slide, new_path):
    for sh in slide.shapes:
        if sh.shape_type == 13:
            blip = sh._element.find('.//' + qn('a:blip'))
            rId = blip.get(qn('r:embed'))
            part = sh.part.related_part(rId)
            part._blob = Path(new_path).read_bytes()
            return True
    return False


def set_textbox(slide, match_sub, new_lines):
    """현재 텍스트에 match_sub 포함된 textbox를 찾아 new_lines로 재작성."""
    for sh in slide.shapes:
        if sh.has_text_frame and match_sub in sh.text_frame.text:
            tf = sh.text_frame
            tf.clear()
            for i, (text, kw) in enumerate(new_lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = kw.get("align", PP_ALIGN.CENTER)
                if "space_before" in kw: p.space_before = Pt(kw["space_before"])
                run = p.add_run(); run.text = text
                _set_font(run, size=kw.get("size", 13), bold=kw.get("bold", False),
                          italic=kw.get("italic", False), color=kw.get("color", NAVY))
            return True
    return False


def set_table_cell(slide, r, c, text):
    for sh in slide.shapes:
        if sh.has_table:
            cell = sh.table.cell(r, c)
            # 첫 run의 서식 유지하며 텍스트만 교체
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = text
                for extra in para.runs[1:]: extra.text = ""
            else:
                para.add_run().text = text
            return


# ── 슬라이드 13: 학습 여정 (chart_milestones) ────────────────────────────────
replace_pic(slides[12], ASSETS / "chart_milestones.png")
set_textbox(slides[12], "데이터 큐레이션",
            [("성능 점프는 모델 구조가 아니라 데이터 큐레이션(+20.6%p)과 모델 업그레이드(+8.0%p)에서",
              dict(size=13, bold=True, color=NAVY))])

# ── 슬라이드 14: 4B 단계별 (chart_progression) ───────────────────────────────
replace_pic(slides[13], ASSETS / "chart_progression.png")
set_textbox(slides[13], "어댑터 없음",
            [("어댑터 없음 0% → Base 14.7% → SFT 85.3% → DPO 85.3% (held-out n=150)",
              dict(size=13, italic=True, color=GRAY))])

# ── 슬라이드 15: 시나리오별 (chart_scenario) ─────────────────────────────────
replace_pic(slides[14], ASSETS / "chart_scenario.png")
set_textbox(slides[14], "의존성 체인만",
            [("대부분 시나리오 90~100% · 의존성 체인만 47%로 공통 막힘 (SFT·DPO 동일, n=150)",
              dict(size=13, bold=True, color=NAVY))])

# ── 슬라이드 17: KEY FINDING ① 표 (n=30 → n=150) ─────────────────────────────
# 표: 행1 instruct 0.0%, 행2 base 16.7%, 행3 +SFT 90.0%
set_table_cell(slides[16], 0, 1, "통과율(n=150)")
set_table_cell(slides[16], 2, 1, "14.7%")   # base
set_table_cell(slides[16], 3, 1, "85.3%")   # +SFT
# 본문: 내용만 채점 시 instruct 34% 추가 설명
set_textbox(slides[16], "instruct 모델 0%",
            [("\"instruct 0%\"는 추론 실패가 아니라 스키마 미준수 — tasks를 [\"문자열\"]로 출력해 파싱 실패.",
              dict(size=13, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)),
             ("내용만 채점(포맷 무시)하면 instruct 34.0% · base 31.3% — 추론 능력의 1/3은 이미 있었고, 파인튜닝이 앱 규격으로 고정한 것.",
              dict(size=11.5, color=DARK, align=PP_ALIGN.LEFT, space_before=4))])

# ── 슬라이드 18: KEY FINDING ② 표 (체인 67 → 47) ─────────────────────────────
set_table_cell(slides[17], 1, 1, "90~100%")  # 비체인 SFT
set_table_cell(slides[17], 2, 1, "47%")       # 체인 SFT
set_table_cell(slides[17], 2, 2, "47%")       # 체인 DPO

# ── 슬라이드 19: 결론 수치 갱신 ───────────────────────────────────────────────
set_textbox(slides[18], "최종 성과",
            [("최종 성과: Qwen3.5-4B 85.3% · Qwen3.5-9B 88.7% (RTX 4090) — 둘 다 held-out n=150",
              dict(size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER))])

# ── 신규 실험 슬라이드 (결론 앞에 삽입) ───────────────────────────────────────
def footer_pg(slide, n):
    add_text(slide, Inches(11.8), Inches(7.05), Inches(1.4), Inches(0.35),
             [(str(n), dict(size=10, color=GRAY, align=PP_ALIGN.RIGHT))])


# 4B vs 9B
s = new_slide("RESULTS", "모델 크기 효과 — Qwen3.5 4B vs 9B")
add_image_fit(s, ASSETS / "chart_4b_vs_9b.png", Inches(0.5), Inches(1.45), Inches(12.35), Inches(4.5))
add_text(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.0),
         [("동일 표본(n=150): 4B 85.3% vs 9B 88.7% (+3.4%p) — 9B는 당일시각·체인에서 앞섬",
           dict(size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)),
          ("→ 체인은 4B 47%·9B 57%로 둘 다 미해결 — 모델 크기가 아니라 SFT 데이터 보강 문제",
           dict(size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER, space_before=3))])

# 9B 4-way
s = new_slide("RESULTS", "Qwen3.5-9B 4-way — 포맷 vs 추론 (n=30)")
add_image_fit(s, ASSETS / "chart_9b_schema_vs_content.png", Inches(0.5), Inches(1.45), Inches(12.35), Inches(4.6))
add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9),
         [("9B도 4B와 동일 패턴 — no-adapter 스키마 0%지만 내용 46.7%(포맷만 미준수), SFT가 93.3%로 고정",
           dict(size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)),
          ("SFT=DPO 완전 동일(체인 4/6 고착) — 9B에서도 DPO가 체인을 못 옮김",
           dict(size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER, space_before=3))])

# 학습 지표
s = new_slide("TRAINING METRICS", "학습 지표 비교 — 4B vs 9B")
add_image_fit(s, ASSETS / "chart_train_metrics_4b_9b.png", Inches(0.4), Inches(1.4), Inches(7.3), Inches(4.4))
add_table(s, Inches(7.9), Inches(1.5), Inches(5.0), Inches(4.6),
          [["지표", "4B", "9B"], ["파라미터", "4.2B", "9.4B"], ["GPU", "3080Ti 12G", "4090 24G"],
           ["SFT train_loss", "0.309", "0.320"], ["SFT token acc", "90.8%", "91.5%"],
           ["SFT 시간", "9.5h", "11.5h"], ["DPO reward_acc", "97.5%", "88.0%"],
           ["DPO margin", "3.50", "0.099"], ["DPO 시간", "~0.5h", "2.5h"],
           ["검증(n=150)", "85.3%", "88.7%"]],
          col_widths=[Inches(2.3), Inches(1.45), Inches(1.25)], font_size=11.5)
add_text(s, Inches(0.4), Inches(6.05), Inches(12.5), Inches(1.0),
         [("SFT(동일 데이터): 9B도 4B와 거의 동일 — 모델 2.3×로도 SFT 학습 이득 미미",
           dict(size=12.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)),
          ("DPO(데이터 다름): margin 차이(3.50 vs 0.099)는 모델이 아닌 데이터 난이도(on-policy vs v4_extra) 탓",
           dict(size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER, space_before=3))])

# ── 신규 3장을 결론(마지막) 앞으로 이동 ───────────────────────────────────────
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
concl = ids[18]            # 원래 19번째(결론) sldId
new_ids = ids[19:]         # 방금 추가된 3장
for nid in new_ids:
    sldIdLst.remove(nid)
    concl.addprevious(nid)

# ── footer 페이지 번호 재부여 (전 슬라이드 순서대로) ─────────────────────────
for i, sid in enumerate(list(prs.slides._sldIdLst), 1):
    sl = prs.slides._sldIdLst  # noqa
# 페이지 번호는 슬라이드 객체로 재계산
for i, sl in enumerate(prs.slides, 1):
    # 기존 footer(숫자만 든 우하단 textbox) 갱신, 없으면 추가
    done = False
    for sh in sl.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t.isdigit() and sh.left and sh.left > Inches(11):
                sh.text_frame.paragraphs[0].runs[0].text = str(i)
                done = True; break
    if not done and i > 1:
        footer_pg(sl, i)

prs.save(OUT)
print(f"[saved] {OUT}")
print(f"총 슬라이드: {len(prs.slides._sldIdLst)}")
