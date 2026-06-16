#!/usr/bin/env python3
"""v5 발표자료 = DROP 앞부분(데이터셋, 사용자 편집본) + v3 뒷부분(실험결과 n=150).

- Base  = DROP deck (사용자가 편집한 앞부분을 byte-perfect 보존)
- 앞부분 = DROP slides [00..15]  (title → PIPELINE)
- 뒷부분 = V3   slides [18..27]  (RESULTS → 결론), 교차파일 복사
- 페이지번호(우하단 숫자 텍스트박스)는 최종 위치 기준 index+1로 재부여
"""
import copy
import io
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

DROP = "/tmp/cmux-drop-9e568278-e565-4b59-9f1b-662597c3a1f7.pptx"
V3 = "presentation/TimeSorter_발표자료_v3.pptx"
OUT = "presentation/TimeSorter_발표자료_v5.pptx"

DROP_KEEP = range(0, 16)   # [00..15] 앞부분 보존
DROP_DROP = range(16, 23)  # [16..22] 옛 뒷부분 삭제
V3_BACK = range(18, 28)    # [18..27] 새 뒷부분

RID_ATTRS = (qn("r:embed"), qn("r:link"), qn("r:id"))

SHAPE_TAGS = (qn("p:sp"), qn("p:pic"), qn("p:graphicFrame"),
              qn("p:grpSp"), qn("p:cxnSp"), qn("p:contentPart"))


def delete_slide(prs, idx):
    """sldIdLst 항목 + prs 관계 둘 다 제거 (고아 part 방지)."""
    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[idx]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def clear_placeholders(slide):
    spTree = slide.shapes._spTree
    for el in list(spTree):
        if el.tag in SHAPE_TAGS:
            spTree.remove(el)


def relink(src_slide, dst_slide, el):
    """복사된 element 내 모든 rId 참조를 dst part로 재연결."""
    for node in el.iter():
        for attr in RID_ATTRS:
            rId = node.get(attr)
            if not rId:
                continue
            rel = src_slide.part.rels[rId]
            if rel.is_external:
                new_rId = dst_slide.part.rels.get_or_add_ext_rel(
                    rel.reltype, rel.target_ref)
            else:
                blob = rel.target_part.blob
                # 이미지면 get_or_add_image_part (partname 충돌 회피)
                if "image" in rel.reltype:
                    _, new_rId = dst_slide.part.get_or_add_image_part(
                        io.BytesIO(blob))
                else:
                    new_rId = dst_slide.part.relate_to(
                        rel.target_part, rel.reltype)
            node.set(attr, new_rId)


def copy_slide(src_slide, dst_prs, blank_layout):
    dst_slide = dst_prs.slides.add_slide(blank_layout)
    clear_placeholders(dst_slide)
    for shape in src_slide.shapes:
        new_el = copy.deepcopy(shape._element)
        dst_slide.shapes._spTree.append(new_el)
        relink(src_slide, dst_slide, new_el)
    return dst_slide


def renumber(prs):
    """우하단 숫자-only 텍스트박스를 index+1로 갱신."""
    changed = []
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            if t.isdigit() and Emu(sh.left).inches > 10.5 and Emu(sh.top).inches > 6.5:
                want = str(i + 1)
                if t != want:
                    # 첫 run에 기록 (서식 보존)
                    p = sh.text_frame.paragraphs[0]
                    if p.runs:
                        p.runs[0].text = want
                        for r in p.runs[1:]:
                            r.text = ""
                    else:
                        p.text = want
                    changed.append((i, t, want))
    return changed


def main():
    blank_name = "Blank"
    prs = Presentation(DROP)
    blank_layout = next(l for l in prs.slide_layouts if l.name == blank_name)

    # 1) DROP 옛 뒷부분 [16..22] 삭제 (뒤에서부터)
    for idx in sorted(DROP_DROP, reverse=True):
        delete_slide(prs, idx)
    assert len(prs.slides._sldIdLst) == len(DROP_KEEP), len(prs.slides._sldIdLst)

    # 2) V3 뒷부분 [18..27] append
    v3 = Presentation(V3)
    v3_slides = list(v3.slides)
    for idx in V3_BACK:
        copy_slide(v3_slides[idx], prs, blank_layout)

    # 3) 페이지번호 재부여
    changed = renumber(prs)

    prs.save(OUT)
    print(f"saved {OUT}: {len(list(prs.slides))} slides")
    print(f"page-number changes ({len(changed)}):")
    for i, old, new in changed:
        print(f"  slide[{i:02d}] {old} -> {new}")


if __name__ == "__main__":
    main()
