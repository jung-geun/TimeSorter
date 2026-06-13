#!/usr/bin/env python3
"""TimeSorter 발표자료 PPT 생성 (python-pptx).

presentation/ 문서 + assets/ 차트를 16:9 슬라이드 덱으로 조립.
출력: presentation/TimeSorter_발표자료.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Inches, Emu

ROOT = Path("/mnt/hdd/WD_8TB/code/TimeSorter")
ASSETS = ROOT / "assets"
OUT = ROOT / "presentation" / "TimeSorter_발표자료.pptx"

# ── 색상 팔레트 ──────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1F, 0x2D, 0x3D)
BLUE    = RGBColor(0x4C, 0x72, 0xB0)
ORANGE  = RGBColor(0xDD, 0x84, 0x52)
GREEN   = RGBColor(0x55, 0xA8, 0x68)
RED     = RGBColor(0xC0, 0x39, 0x2B)
GRAY    = RGBColor(0x5A, 0x5A, 0x5A)
LIGHT   = RGBColor(0xF4, 0xF6, 0xF8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x22, 0x22, 0x22)

FONT = "Noto Sans CJK KR"      # 뷰어에 없으면 Malgun Gothic 등으로 대체됨
MONO = "D2Coding"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set_font(run, size=18, bold=False, color=DARK, font=FONT, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font


def add_box(slide, x, y, w, h, fill=None, line=None, line_w=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    return shp


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True):
    """lines: list of (text, dict) — dict는 _set_font kwargs + optional 'bullet','space_after','level'."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for i, (text, kw) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = kw.pop("align", align)
        if "space_after" in kw:
            p.space_after = Pt(kw.pop("space_after"))
        if "space_before" in kw:
            p.space_before = Pt(kw.pop("space_before"))
        p.level = kw.pop("level", 0)
        bullet = kw.pop("bullet", False)
        run = p.add_run()
        run.text = ("•  " if bullet else "") + text
        _set_font(run, **kw)
    return tb


def title_bar(slide, kicker, title):
    add_box(slide, 0, 0, SW, Inches(1.15), fill=NAVY)
    add_box(slide, 0, Inches(1.15), SW, Pt(3), fill=ORANGE)
    if kicker:
        add_text(slide, Inches(0.5), Inches(0.14), Inches(12), Inches(0.35),
                 [(kicker, dict(size=12, bold=True, color=ORANGE))])
    add_text(slide, Inches(0.5), Inches(0.42), Inches(12.3), Inches(0.65),
             [(title, dict(size=26, bold=True, color=WHITE))])


def content_slide(kicker, title):
    s = prs.slides.add_slide(BLANK)
    title_bar(s, kicker, title)
    return s


def footer(slide, n):
    add_text(slide, Inches(11.8), Inches(7.05), Inches(1.4), Inches(0.35),
             [(f"{n}", dict(size=10, color=GRAY, align=PP_ALIGN.RIGHT))])


# ── 표 헬퍼 ──────────────────────────────────────────────────────────────────
def add_table(slide, x, y, w, h, data, col_widths=None, header_fill=BLUE,
              font_size=13, header_size=13):
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            gt.columns[ci].width = cw
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.add_run(); run.text = str(val)
            if ri == 0:
                _set_font(run, size=header_size, bold=True, color=WHITE)
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                _set_font(run, size=font_size, color=DARK)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
    return gt


def add_image_fit(slide, path, x, y, max_w, max_h):
    from PIL import Image
    iw, ih = Image.open(path).size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    slide.shapes.add_picture(str(path), x + (max_w - w) // 2, y + (max_h - h) // 2, w, h)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_box(s, 0, 0, SW, SH, fill=NAVY)
add_box(s, 0, Inches(4.35), SW, Pt(4), fill=ORANGE)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.3),
         [("TimeSorter", dict(size=54, bold=True, color=WHITE))])
add_text(s, Inches(0.95), Inches(3.5), Inches(11.5), Inches(0.8),
         [("한국어 일정 우선순위 정렬 비서 — Qwen3.5-4B 파인튜닝", dict(size=22, color=RGBColor(0xC9,0xD4,0xE0)))])
add_text(s, Inches(0.95), Inches(4.6), Inches(11.5), Inches(1.2),
         [("SFT → DPO 파이프라인 · 4축 JSON 채점 · RTX 12GB QLoRA",
           dict(size=15, color=ORANGE)),
          ("할 일 목록을 긴급도·중요도·의존성·시간제약 4축으로 채점해 실행 순서를 제안",
           dict(size=13, color=RGBColor(0x9F,0xAE,0xBF), space_before=8))])
add_text(s, Inches(0.95), Inches(6.7), Inches(11), Inches(0.4),
         [("2026-06 · github.com/jung-geun/TimeSorter", dict(size=11, color=GRAY))])

# ═══════════════════════════════════════════════════════════════════════════
# Slide 2 — 프로젝트 개요
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("OVERVIEW", "프로젝트 개요 — 무엇을 푸는가")
add_text(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.9),
         [("스마트폰·PC에서 '오늘 할 일'을 입력하면, AI가 페르소나와 맥락을 이해해 "
           "실행 순서를 제안하는 개인 비서 코어 모델.",
           dict(size=16, color=DARK)),
          ("단순 키워드 정렬이 아니라, 각 태스크를 4개 축으로 1–5점 채점하고 근거를 함께 제시한다.",
           dict(size=14, color=GRAY, space_before=6))])

add_box(s, Inches(0.5), Inches(3.0), Inches(5.9), Inches(3.7), fill=LIGHT, line=BLUE, line_w=Pt(1.5))
add_text(s, Inches(0.7), Inches(3.15), Inches(5.5), Inches(0.4),
         [("입력", dict(size=15, bold=True, color=BLUE))])
add_text(s, Inches(0.7), Inches(3.6), Inches(5.6), Inches(3.0),
         [("[페르소나] 최건호 (플랜트공학 기술자, 63세)", dict(size=12, bold=True, color=DARK, space_after=6)),
          ("오늘: 2026-03-30", dict(size=11, color=GRAY, space_after=8)),
          ("- 배관 열손실 검토표 제출 (3/30 09:30)", dict(size=12, color=DARK)),
          ("- 공정안전성 회의자료 초안 (4/4 10:00)", dict(size=12, color=DARK)),
          ("- 현장 펌프 진동 기록 정리 (마감 없음)", dict(size=12, color=DARK)),
          ("- 열교환기 시험 보고서 (3/30 15:00)", dict(size=12, color=DARK)),
          ("- 부식 코팅 검토 의견 (3/29 18:00 → 지남)", dict(size=12, color=RED))])

add_box(s, Inches(6.9), Inches(3.0), Inches(5.9), Inches(3.7), fill=LIGHT, line=GREEN, line_w=Pt(1.5))
add_text(s, Inches(7.1), Inches(3.15), Inches(5.5), Inches(0.4),
         [("출력 (4축 JSON)", dict(size=15, bold=True, color=GREEN))])
add_text(s, Inches(7.1), Inches(3.6), Inches(5.6), Inches(3.0),
         [("우선순위: 1 → 4 → 2 → 3 → 5", dict(size=13, bold=True, color=DARK, space_after=6)),
          ("① 배관 검토표  [긴급5·중요4·의존1·시간5]", dict(size=11.5, color=DARK)),
          ("④ 열교환기 보고서 [긴급4·중요4·의존2·시간4]", dict(size=11.5, color=DARK)),
          ("② 회의자료 초안 [긴급2·중요4·의존3·시간2]", dict(size=11.5, color=DARK)),
          ("③ 펌프 기록 정리 [긴급1·중요2·의존1·시간1]", dict(size=11.5, color=DARK)),
          ("⑤ 부식 검토 [긴급1·중요1·의존1·시간1] 이미 지남", dict(size=11.5, color=RED, space_before=2)),
          ("→ 지난 일정은 자동 강등·최하위 배치", dict(size=11, italic=True, color=GREEN, space_before=6))])
footer(s, 2)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 3 — 데이터 출처
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("DATA SOURCE", "원본 데이터셋과 증강 — 이식에서 합성으로")
add_box(s, Inches(0.5), Inches(1.45), Inches(6.0), Inches(2.55), fill=LIGHT, line=BLUE)
add_text(s, Inches(0.7), Inches(1.55), Inches(5.6), Inches(2.4),
         [("① anakin87/events-scheduling (영문, 500행)", dict(size=14, bold=True, color=BLUE, space_after=4)),
          ("원본: 시간 겹치는 이벤트에서 충돌 없이 점수 최대화하는 구간 스케줄링 최적화", dict(size=11.5, color=DARK, space_after=4)),
          ("변형: 영문 이벤트 → 한국어 할 일로 번역, '최적화'를 '4축 우선순위 채점'으로 재정의", dict(size=11.5, color=DARK)),
          ("(예: Small Language Models talk → 작은 언어 모델 이야기 듣기)", dict(size=10.5, italic=True, color=GRAY, space_before=3))])
add_box(s, Inches(6.85), Inches(1.45), Inches(6.0), Inches(2.55), fill=LIGHT, line=GREEN)
add_text(s, Inches(7.05), Inches(1.55), Inches(5.6), Inches(2.4),
         [("② nvidia/Nemotron-Personas-Korea (한국어, 100만행)", dict(size=14, bold=True, color=GREEN, space_after=4)),
          ("원본: 직업·취미·가족 등 다축으로 묘사된 합성 한국어 페르소나", dict(size=11.5, color=DARK, space_after=4)),
          ("변형: '이름 (직업, 나이)'로 압축해 시스템 프롬프트에 주입", dict(size=11.5, color=DARK)),
          ("→ 같은 할 일도 페르소나마다 다르게 정렬 (범용 정렬기가 아닌 개인 비서)", dict(size=10.5, italic=True, color=GRAY, space_before=3))])

add_box(s, Inches(0.5), Inches(4.3), Inches(12.35), Inches(2.4), fill=RGBColor(0xFD,0xF3,0xE7), line=ORANGE)
add_text(s, Inches(0.7), Inches(4.45), Inches(12), Inches(2.2),
         [("왜 그대로 쓰지 않았나", dict(size=14, bold=True, color=ORANGE, space_after=5)),
          ("anakin87 이벤트는 날짜·마감·선후관계가 없는 추상적 강연 목록 → 우리 과제(날짜 추론·의존성·리스크)에 부족",
           dict(size=12.5, color=DARK, bullet=True)),
          ("v3부터 골격(skeleton) 생성기로 현실적 한국어 일정을 합성하며 원본 이벤트 의존을 점진적으로 제거",
           dict(size=12.5, color=DARK, bullet=True)),
          ("Nemotron 페르소나는 끝까지 조건부 신호로 유지 → '이식(graft) → 합성(synthesis)' 전략",
           dict(size=12.5, color=DARK, bullet=True))])
footer(s, 3)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 4 — 출력 포맷 규격
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("OUTPUT FORMAT", "출력 포맷 규격 — 왜 4축 JSON인가")
add_box(s, Inches(0.5), Inches(1.5), Inches(6.3), Inches(5.1), fill=RGBColor(0x2B,0x2B,0x2B))
add_text(s, Inches(0.7), Inches(1.62), Inches(6.0), Inches(4.9),
         [('{', dict(size=12, color=WHITE, font=MONO)),
          ('  "tasks": [{"id":1, "text":"원문 그대로"}],', dict(size=12, color=RGBColor(0x9C,0xDC,0xFE), font=MONO)),
          ('  "priority_order": [2, 3, 1],', dict(size=12, color=RGBColor(0xCE,0x91,0x78), font=MONO)),
          ('  "scores": [', dict(size=12, color=WHITE, font=MONO)),
          ('    {"task_id":1, "urgency":3,', dict(size=12, color=RGBColor(0xB5,0xCE,0xA8), font=MONO)),
          ('     "importance":2, "dependency":1,', dict(size=12, color=RGBColor(0xB5,0xCE,0xA8), font=MONO)),
          ('     "time_constraint":1,', dict(size=12, color=RGBColor(0xB5,0xCE,0xA8), font=MONO)),
          ('     "reason":"한 문장 근거"}', dict(size=12, color=RGBColor(0xB5,0xCE,0xA8), font=MONO)),
          ('  ],', dict(size=12, color=WHITE, font=MONO)),
          ('  "refusal_reason": ""', dict(size=12, color=RGBColor(0xCE,0x91,0x78), font=MONO)),
          ('}', dict(size=12, color=WHITE, font=MONO))])
add_text(s, Inches(7.0), Inches(1.5), Inches(5.9), Inches(5.1),
         [("규격화 의도 4가지", dict(size=16, bold=True, color=NAVY, space_after=8)),
          ("JSON 단일 출력 강제", dict(size=13.5, bold=True, color=BLUE, bullet=True)),
          ("앱이 후처리 없이 바로 파싱·렌더링. 자유 텍스트는 연동 불가", dict(size=11.5, color=GRAY, level=1, space_after=6)),
          ("점수와 순서의 분리", dict(size=13.5, bold=True, color=BLUE, bullet=True)),
          ("scores(왜) ↔ priority_order(결과) 분리 → 후처리 ScoreRanker로 모순 교정", dict(size=11.5, color=GRAY, level=1, space_after=6)),
          ("4축 다축 채점", dict(size=13.5, bold=True, color=BLUE, bullet=True)),
          ("단순 긴급순이 아닌 긴급·중요·의존·시간을 독립 채점", dict(size=11.5, color=GRAY, level=1, space_after=6)),
          ("원문 보존 (tasks.text)", dict(size=13.5, bold=True, color=BLUE, bullet=True)),
          ("id-텍스트 매핑 강제로 항목 요약·누락 차단", dict(size=11.5, color=GRAY, level=1, space_after=6)),
          ("→ 골격 규칙으로 $0·결정론적 자동 채점 가능", dict(size=12, bold=True, italic=True, color=GREEN, space_before=4))])
footer(s, 4)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 5 — 4축 채점 기준
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("SCORING", "4축 채점 기준 (각 1–5점)")
add_table(s, Inches(0.5), Inches(1.6), Inches(12.35), Inches(3.6),
          [["축", "의미", "5점", "1점"],
           ["urgency\n긴급도", "마감/시작 임박도", "오늘 마감 (오전=5·오후=4)", "마감 없음 / 이미 지난 일정"],
           ["importance\n중요도", "페르소나 목표 영향도", "결정적 (에스컬레이션·위약금·법적기한)", "옵션 (없어도 무방)"],
           ["dependency\n의존성", "후속 작업 블로킹 정도", "다수 후속의 선행 입력", "완전 독립"],
           ["time_constraint\n시간제약", "고정 시각 강도", "시각 고정 (회의 등)", "언제든 가능"]],
          col_widths=[Inches(2.3), Inches(3.2), Inches(4.0), Inches(2.85)],
          font_size=13)
add_box(s, Inches(0.5), Inches(5.45), Inches(12.35), Inches(1.3), fill=RGBColor(0xFD,0xF3,0xE7), line=ORANGE)
add_text(s, Inches(0.7), Inches(5.58), Inches(12), Inches(1.1),
         [("우선순위 결정 로직 (단순 합산 아님)", dict(size=13, bold=True, color=ORANGE, space_after=3)),
          ("① 이미 지난 일정 → 무조건 최하위(urgency·tc 강제=1)   ② 의존성 체인 → 선행부터 연속 배치   "
           "③ 나머지 → 4축 종합 판단", dict(size=12, color=DARK))])
footer(s, 5)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 6 — SFT vs DPO 데이터셋
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("DATASET", "SFT vs DPO 데이터셋 개형")
add_box(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.7), fill=LIGHT, line=BLUE)
add_text(s, Inches(0.7), Inches(1.6), Inches(5.6), Inches(2.5),
         [("SFT — 지도학습 (sft_v4_train, 6,056행)", dict(size=14, bold=True, color=BLUE, space_after=4)),
          ("'이 입력 → 이 정답' 쌍만 필요", dict(size=12, color=DARK, space_after=4)),
          ("prompt · chosen(정답 JSON) · persona · today · source", dict(size=11.5, color=GRAY)),
          ("목적: 4축 JSON 포맷과 채점 규칙 자체를 주입 (없는 능력 추가)", dict(size=11.5, italic=True, color=BLUE, space_before=6))])
add_box(s, Inches(6.85), Inches(1.5), Inches(6.0), Inches(2.7), fill=LIGHT, line=ORANGE)
add_text(s, Inches(7.05), Inches(1.6), Inches(5.6), Inches(2.5),
         [("DPO — 선호학습 (dpo_pairs_v5, 1,368쌍)", dict(size=14, bold=True, color=ORANGE, space_after=4)),
          ("'좋은 답(chosen) vs 나쁜 답(rejected)' 쌍 필요", dict(size=12, color=DARK, space_after=4)),
          ("+ rejected · category(위반유형)", dict(size=11.5, color=GRAY)),
          ("목적: 한 가지 규칙만 틀린 hard negative로 그 규칙 위반을 벌점", dict(size=11.5, italic=True, color=ORANGE, space_before=6))])
add_text(s, Inches(0.5), Inches(4.35), Inches(12), Inches(0.4),
         [("DPO rejected 위반 유형 분포 (v5)", dict(size=13, bold=True, color=NAVY))])
add_table(s, Inches(0.5), Inches(4.8), Inches(12.35), Inches(1.9),
          [["유형", "쌍", "rejected가 일부러 저지른 오류"],
           ["order_score_mismatch", "230", "점수는 맞는데 순서가 점수와 모순"],
           ["granularity_swap", "220", "마감 단위(시각↔날짜) 혼동"],
           ["date_confusion", "177", "지난 날짜를 오늘/미래로 착각"],
           ["dependency_scatter", "185", "의존성 체인을 분산 배치"],
           ["risk_ignore / past_hallucination", "92 / 36", "리스크 미반영 / 지남 단정"]],
          col_widths=[Inches(4.3), Inches(1.8), Inches(6.25)], font_size=12, header_fill=ORANGE)
footer(s, 6)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 7 — DPO 선호쌍 예시
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("DPO PAIR", "DPO 선호쌍 실제 예시 — date_confusion")
add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.0),
         [("입력: 정순자 (무직, 86세) · 오늘 = 2026-12-11", dict(size=13, bold=True, color=DARK, space_after=3)),
          ("id=4 '12월 9일 14:00까지 손주 겨울 인사 카드 써두기' → 오늘(12/11)보다 2일 지난 일정",
           dict(size=12.5, color=RED)),
          ("두 응답의 차이는 오직 이 한 태스크(id=4)뿐", dict(size=11.5, italic=True, color=GRAY, space_before=2))])
add_box(s, Inches(0.5), Inches(3.0), Inches(6.0), Inches(3.5), fill=RGBColor(0xEC,0xF6,0xEE), line=GREEN, line_w=Pt(2))
add_text(s, Inches(0.7), Inches(3.15), Inches(5.6), Inches(3.2),
         [("✅ chosen (선호)", dict(size=15, bold=True, color=GREEN, space_after=6)),
          ("priority_order: [1, 3, 5, 2, 4]", dict(size=13, bold=True, color=DARK, font=MONO)),
          ("→ id=4를 최하위(5번째) 배치", dict(size=12, color=DARK, space_after=8)),
          ("id=4 점수: urgency=1, time_constraint=1", dict(size=12.5, color=DARK, font=MONO)),
          ('reason: "12월 9일 마감으로 오늘 기준', dict(size=11.5, italic=True, color=GREEN, space_before=4)),
          (' 이미 지난 일정이므로 최하위로 배치"', dict(size=11.5, italic=True, color=GREEN))])
add_box(s, Inches(6.85), Inches(3.0), Inches(6.0), Inches(3.5), fill=RGBColor(0xFB,0xEC,0xEA), line=RED, line_w=Pt(2))
add_text(s, Inches(7.05), Inches(3.15), Inches(5.6), Inches(3.2),
         [("❌ rejected (비선호)", dict(size=15, bold=True, color=RED, space_after=6)),
          ("priority_order: [4, 1, 3, 5, 2]", dict(size=13, bold=True, color=DARK, font=MONO)),
          ("→ id=4를 1위 배치", dict(size=12, color=DARK, space_after=8)),
          ("id=4 점수: urgency=5, time_constraint=5", dict(size=12.5, color=DARK, font=MONO)),
          ('reason: "마감이 오늘이라', dict(size=11.5, italic=True, color=RED, space_before=4)),
          (' 즉시 착수해야 합니다" (12/9를 오늘로 착각)', dict(size=11.5, italic=True, color=RED))])
add_text(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.6),
         [("포맷은 100% 동일, 내용 한 곳만 다름 → 학습 신호가 '날짜 추론 규칙'에 정확히 집중 (hard negative)",
           dict(size=12, bold=True, italic=True, color=NAVY, align=PP_ALIGN.CENTER))])
footer(s, 7)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 8 — 데이터셋 진화
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("EVOLUTION", "데이터셋 진화 v1 → v5")
add_table(s, Inches(0.5), Inches(1.55), Inches(12.35), Inches(3.2),
          [["버전", "출력 형식", "핵심 추가", "기대효과"],
           ["v1", "자유 텍스트", "페르소나 기반 우선순위", "말투·관점 반영 (앱 파싱 불가)"],
           ["v2", "4축 점수 JSON", "구조화·앱연동·설명가능성", "정렬 근거 정량화 (날짜 개념 없음)"],
           ["v3", "4축 JSON", "오늘 날짜 주입·지난 일정 규칙", "날짜 혼동 해결 (dated 38→77%)"],
           ["v4", "4축 JSON", "이중 생성·prompt loss 마스킹", "56.7 → 77.3% (+20.6%p)"],
           ["v5", "4축 JSON", "on-policy DPO·체인 보강", "성능 보존 + 날짜 미세 개선"]],
          col_widths=[Inches(1.1), Inches(2.6), Inches(4.2), Inches(4.45)], font_size=12.5)
add_box(s, Inches(0.5), Inches(5.0), Inches(12.35), Inches(1.7), fill=LIGHT, line=BLUE)
add_text(s, Inches(0.7), Inches(5.13), Inches(12), Inches(1.5),
         [("핵심 흐름", dict(size=13, bold=True, color=BLUE, space_after=4)),
          ("자유 텍스트(v1) → 구조화 JSON(v2) → 날짜 추론(v3) → 데이터 품질·이중생성(v4) → on-policy 보강(v5)",
           dict(size=12.5, color=DARK, space_after=3)),
          ("성능 점프(+20.6%p)는 모델 구조가 아니라 데이터 큐레이션 + prompt loss 마스킹에서 나왔다",
           dict(size=12, bold=True, italic=True, color=ORANGE))])
footer(s, 8)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 9 — 학습 파이프라인
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("PIPELINE", "학습 파이프라인 & RTX 12GB 제약 극복")
# 파이프라인 흐름
steps = [("Base\nQwen3.5-4B", GRAY), ("SFT v4\n포맷·규칙 주입", BLUE),
         ("DPO v5\n선호 미세조정", ORANGE)]
x = Inches(0.7)
for i, (label, color) in enumerate(steps):
    add_box(s, x, Inches(1.7), Inches(2.6), Inches(1.1), fill=color)
    add_text(s, x, Inches(1.83), Inches(2.6), Inches(0.95),
             [(label, dict(size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER))],
             anchor=MSO_ANCHOR.MIDDLE)
    if i < 2:
        add_text(s, x + Inches(2.6), Inches(1.83), Inches(0.6), Inches(0.95),
                 [("→", dict(size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER))],
                 anchor=MSO_ANCHOR.MIDDLE)
    x = x + Inches(3.2)
add_text(s, Inches(0.5), Inches(3.1), Inches(12), Inches(0.4),
         [("12GB VRAM 한계 극복 기법", dict(size=15, bold=True, color=NAVY))])
add_text(s, Inches(0.7), Inches(3.6), Inches(12), Inches(3.0),
         [("QLoRA 4-bit NF4 — 4.2B 모델 + 248K vocab을 12GB에 적재", dict(size=13, color=DARK, bullet=True, space_after=5)),
          ("MemEfficientDPOTrainer — chosen/rejected [1,T] 순차 처리 + 청킹 log_softmax + 수동 gradient 주입",
           dict(size=13, color=DARK, bullet=True, space_after=5)),
          ("(248K vocab → [2,T,248K] logits ~1.1GiB, 표준 DPO 불가 → 자체 구현)", dict(size=11, italic=True, color=GRAY, level=1, space_after=5)),
          ("prompt loss 마스킹 — 응답 토큰에만 loss → 학습 밀도 ~2.5배", dict(size=13, color=DARK, bullet=True, space_after=5)),
          ("Qwen3.5 CUDA 버그 우회 — is_fast_path_available 강제 패치 (fla delta-rule)", dict(size=13, color=DARK, bullet=True, space_after=5)),
          ("precompute_ref_log_probs — ref 모델 forward 사전 계산으로 메모리 절감", dict(size=13, color=DARK, bullet=True))])
footer(s, 9)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 10 — 학습 커브
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("RESULTS", "벤치마크 — 학습 여정의 두 점프")
add_image_fit(s, ASSETS / "chart_milestones.png", Inches(1.6), Inches(1.45),
              Inches(10.1), Inches(4.7))
add_text(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.7),
         [("성능 점프는 모델 구조가 아니라 데이터 큐레이션(+20.6%p)과 모델 업그레이드(+12.7%p)에서",
           dict(size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER))])
footer(s, 10)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 11 — Qwen3.5 4모델 비교
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("RESULTS", "Qwen3.5-4B 파인튜닝 단계별 통과율")
add_image_fit(s, ASSETS / "chart_progression.png", Inches(1.6), Inches(1.45),
              Inches(10.1), Inches(4.7))
add_text(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.6),
         [("어댑터 없음 0% → Base 16.7% → SFT 90% → DPO 90% (held-out n=30)",
           dict(size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER))])
footer(s, 11)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 12 — 시나리오별 히트맵
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("RESULTS", "시나리오별 통과율 — SFT vs DPO")
add_image_fit(s, ASSETS / "chart_scenario.png", Inches(1.3), Inches(1.45),
              Inches(10.7), Inches(4.7))
add_text(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.6),
         [("대부분 시나리오 90~100% · 의존성 체인만 67%로 공통 막힘 (SFT·DPO 동일)",
           dict(size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER))])
footer(s, 12)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 13 — 4모델 실제 출력 비교
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("COMPARISON", "4모델 실제 출력 비교 (동일 쿼리)")
add_table(s, Inches(0.5), Inches(1.55), Inches(12.35), Inches(1.5),
          [["쿼리(시나리오)", "Base", "Instruct", "SFT v4", "DPO v5"],
           ["날짜혼재 / 당일시각 / 리스크 / 상대날짜", "—", "—", "—", "—"],
           ["6쿼리 통과 합계", "2/6", "0/6", "5/6", "5/6"]],
          col_widths=[Inches(5.15), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8)],
          font_size=13)
add_box(s, Inches(0.5), Inches(3.4), Inches(6.0), Inches(3.3), fill=RGBColor(0xFB,0xEC,0xEA), line=RED)
add_text(s, Inches(0.65), Inches(3.5), Inches(5.7), Inches(3.1),
         [("Instruct (어댑터 없음) — ❌ 파싱 실패", dict(size=13, bold=True, color=RED, space_after=4)),
          ('"tasks": [', dict(size=11.5, color=DARK, font=MONO)),
          ('  "배관 열손실 검토표를 정리해...",', dict(size=11.5, color=RED, font=MONO)),
          ('  "4월 4일 10:00까지 회의자료...",', dict(size=11.5, color=RED, font=MONO)),
          ('  ... ]', dict(size=11.5, color=DARK, font=MONO)),
          ("tasks를 문자열 배열로 출력 (객체 아님)", dict(size=11, italic=True, color=RED, space_before=5)),
          ("→ 추론은 합리적이나 스키마 위반으로 파싱 실패", dict(size=11, color=GRAY, space_before=3))])
add_box(s, Inches(6.85), Inches(3.4), Inches(6.0), Inches(3.3), fill=RGBColor(0xEC,0xF6,0xEE), line=GREEN)
add_text(s, Inches(7.0), Inches(3.5), Inches(5.7), Inches(3.1),
         [("SFT v4 — ✅ 통과", dict(size=13, bold=True, color=GREEN, space_after=4)),
          ('{"tasks":[{"id":1,"text":"배관..."},', dict(size=11.5, color=GREEN, font=MONO)),
          ('  {"id":2,"text":"회의자료..."}, ...],', dict(size=11.5, color=GREEN, font=MONO)),
          ('"priority_order":[1,4,2,3,5],', dict(size=11.5, color=DARK, font=MONO)),
          ('"scores":[...]}', dict(size=11.5, color=DARK, font=MONO)),
          ("정확한 {id,text} 객체 스키마 준수", dict(size=11, italic=True, color=GREEN, space_before=5)),
          ("→ 지난 일정(id=5) 최하위, urgency=1 강등", dict(size=11, color=GRAY, space_before=3))])
footer(s, 13)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 14 — 핵심 발견 ①
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("KEY FINDING ①", "파인튜닝의 1차 효과는 '스키마 준수'")
add_image_fit(s, ASSETS / "chart_schema_vs_content.png", Inches(0.4), Inches(1.4),
              Inches(8.0), Inches(5.5))
add_box(s, Inches(8.6), Inches(1.5), Inches(4.3), Inches(5.0), fill=LIGHT, line=ORANGE, line_w=Pt(1.5))
add_text(s, Inches(8.78), Inches(1.65), Inches(3.95), Inches(4.8),
         [("\"0%\"의 진실", dict(size=15, bold=True, color=ORANGE, space_after=6)),
          ("no-FT 모델은 tasks를 [\"문자열\"]로 출력해 파싱 실패 → 스키마 0%", dict(size=11.5, color=DARK, bullet=True, space_after=5)),
          ("하지만 priority_order·scores는 정상 → id 재구성 후 내용만 채점하면:", dict(size=11.5, color=DARK, bullet=True, space_after=5)),
          ("instruct 0% → 내용 43.3%", dict(size=12.5, bold=True, color=BLUE, level=1, space_after=2)),
          ("base 16.7% → 내용 40.0%", dict(size=12.5, bold=True, color=BLUE, level=1, space_after=6)),
          ("추론 능력의 40%+ 는 이미 있었다 — 파인튜닝이 그것을 앱 규격으로 고정", dict(size=11.5, color=GRAY, bullet=True, space_after=5)),
          ("FT 모델은 스키마=내용 (갭 없음)", dict(size=11.5, color=GRAY, bullet=True))])
footer(s, 14)

# ─────────────────────────────────────────────────────────────────────────
# Slide 15 — 핵심 발견 ① 보강: 시나리오별 포맷 vs 내용
# ─────────────────────────────────────────────────────────────────────────
s = content_slide("KEY FINDING ①-b", "no-FT 모델의 숨은 추론 능력 — 시나리오별")
add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.55),
         [("Qwen3.5-4B instruct (어댑터 없음) — 스키마 기준 전부 0%지만, 내용만 보면:",
           dict(size=14, bold=True, color=NAVY))])
add_table(s, Inches(0.5), Inches(2.05), Inches(12.35), Inches(2.7),
          [["시나리오", "스키마 통과", "내용 통과", "해석"],
           ["리스크", "0/5", "5/5 (100%)", "추론 완벽 — 포맷만 못 맞춤"],
           ["상대 날짜", "0/4", "4/4 (100%)", "추론 완벽 — 포맷만 못 맞춤"],
           ["당일 시각", "0/6", "2/6 (33%)", "시각 순서 추론 약함"],
           ["날짜 혼재", "0/9", "2/9 (22%)", "지난 일정 처리 약함 (진짜 약점)"],
           ["의존성 체인", "0/6", "0/6 (0%)", "체인 능력 부재 (전 모델 공통)"]],
          col_widths=[Inches(2.6), Inches(2.5), Inches(2.6), Inches(4.65)], font_size=12.5)
add_box(s, Inches(0.5), Inches(5.05), Inches(12.35), Inches(1.7), fill=LIGHT, line=BLUE)
add_text(s, Inches(0.7), Inches(5.18), Inches(12), Inches(1.5),
         [("파인튜닝이 한 일 = 두 종류", dict(size=13, bold=True, color=BLUE, space_after=4)),
          ("① 포맷 고정 (즉효) — 리스크·상대날짜는 내용 이미 100%, SFT는 규격만 고정해 0→100%",
           dict(size=12, color=DARK, bullet=True, space_after=3)),
          ("② 진짜 추론 보강 — 날짜혼재(22→89%)·당일시각(33→100%)은 SFT가 실제 추론을 가르침",
           dict(size=12, color=DARK, bullet=True))])
footer(s, 15)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 16 — 핵심 발견 ② + 위반 추이
# ═══════════════════════════════════════════════════════════════════════════
s = content_slide("KEY FINDING ②", "DPO는 '능력 부재'를 해결하지 못한다")
add_box(s, Inches(0.5), Inches(1.55), Inches(12.35), Inches(1.35), fill=RGBColor(0xFD,0xF3,0xE7), line=ORANGE, line_w=Pt(2))
add_text(s, Inches(0.7), Inches(1.68), Inches(12), Inches(1.15),
         [("의존성 체인에서 SFT v4와 DPO v5가 글자 단위로 동일하게 실패", dict(size=16, bold=True, color=ORANGE, space_after=4)),
          ("같은 priority_order [5,3,4,1,2], 같은 dependency=2 → 동일한 위반 2건. DPO가 체인을 전혀 안 옮겼다.",
           dict(size=12.5, color=DARK))])
add_table(s, Inches(0.5), Inches(3.2), Inches(12.35), Inches(1.5),
          [["시나리오", "SFT v4", "DPO v5", "변화"],
           ["날짜혼재·당일시각·상대날짜", "89~100%", "100%", "유지/소폭 개선"],
           ["의존성 체인", "67%", "67%", "0 (불변)"]],
          col_widths=[Inches(5.35), Inches(2.3), Inches(2.3), Inches(2.4)],
          font_size=13, header_fill=ORANGE)
add_text(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.8),
         [("왜 DPO가 체인을 못 고쳤나", dict(size=14, bold=True, color=NAVY, space_after=5)),
          ("DPO는 '모델이 정답을 생성할 수 있는데 잘못 선택'할 때만 효과적 — 체인은 4B 모델의 능력 부재라 선호학습으로 안 됨",
           dict(size=12.5, color=DARK, bullet=True, space_after=4)),
          ("reward_accuracy 98.9% → 이미 아는 것을 재학습한 셈 (가중치 변화 미미)",
           dict(size=12.5, color=DARK, bullet=True, space_after=4)),
          ("해결책: 4–5단계 긴 체인 특화 SFT 데이터 3× 증량 (선호학습이 아닌 능력 보강)",
           dict(size=12.5, bold=True, color=GREEN, bullet=True))])
footer(s, 16)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 17 — 결론
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_box(s, 0, 0, SW, SH, fill=NAVY)
add_box(s, 0, Inches(1.0), SW, Pt(3), fill=ORANGE)
add_text(s, Inches(0.7), Inches(0.35), Inches(12), Inches(0.7),
         [("결론 & 향후 과제", dict(size=30, bold=True, color=WHITE))])
add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(3.2),
         [("성능 점프는 데이터 품질·loss 마스킹에서 나왔다", dict(size=17, bold=True, color=ORANGE, bullet=True, space_after=4)),
          ("DPO/GRPO 선호학습 단독으로는 한계 약점을 못 옮김 (+20.6%p는 SFT 데이터 큐레이션)", dict(size=13, color=RGBColor(0xC9,0xD4,0xE0), level=1, space_after=10)),
          ("파인튜닝의 1차 효과는 정확한 JSON 스키마 준수", dict(size=17, bold=True, color=ORANGE, bullet=True, space_after=4)),
          ("instruct 0% → SFT 90%: 추론이 아닌 출력 규격 고정이 핵심", dict(size=13, color=RGBColor(0xC9,0xD4,0xE0), level=1, space_after=10)),
          ("의존성 체인은 미해결 — SFT 데이터 증량이 다음 레버", dict(size=17, bold=True, color=ORANGE, bullet=True, space_after=4)),
          ("4–5단계 긴 체인 특화 데이터 3× 증량으로 능력 자체를 보강해야", dict(size=13, color=RGBColor(0xC9,0xD4,0xE0), level=1))])
add_box(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.1), fill=RGBColor(0x2A,0x3B,0x4F), line=ORANGE)
add_text(s, Inches(1.1), Inches(5.75), Inches(11.1), Inches(0.9),
         [("최종 성과: Qwen3.5-4B SFT+DPO = 90.0% (RTX 12GB QLoRA, held-out 자동 채점)",
           dict(size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER))],
         anchor=MSO_ANCHOR.MIDDLE)

prs.save(str(OUT))
print(f"[saved] {OUT}")
print(f"슬라이드 수: {len(prs.slides.__iter__.__self__._sldIdLst)}")
